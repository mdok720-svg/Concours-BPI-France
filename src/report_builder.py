"""
report_builder.py
=================

Assemble le rapport PowerPoint final à partir de :
- presentation.pptx comme TEMPLATE (charte + logos + polices Bpifrance)
- FullReport (données et métriques, cf. indicators.py)
- FullCommentary (textes générés par Gemini, cf. commentary.py)
- Dict des chemins PNG des graphiques (cf. charts.py)

Stratégie :
- Ouvrir le template avec python-pptx
- Identifier les "slots" à remplacer par leur shape_id (IDs stables
  dans le XML du template)
- Remplacer les textes en préservant le formatting du 1er run original
- Supprimer les CHART natifs et insérer nos PNG à la même position
- Sauvegarder dans output/rapport_conjoncture.pptx

Pourquoi le shape_id plutôt que le nom ou la position :
- Les IDs sont stables : si Bpifrance renomme un placeholder ou déplace
  une zone de quelques pixels dans le template, on ne casse pas le code.
- Les noms peuvent être dupliqués ("Espace réservé du contenu 1" apparaît
  plusieurs fois dans le template).
- Les positions peuvent bouger sans préavis.

Auteur : projet challenge Bpifrance — phase 4 (assemblage PPTX).
"""

from __future__ import annotations

import re
from copy import deepcopy
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor

from .indicators import FullReport
from .commentary import FullCommentary


# ---------------------------------------------------------------------------
# Mapping des shape_id du template
#
# Extraits par inspection de presentation.pptx (script phase 4). Si le
# template Bpifrance est un jour remplacé, c'est la seule table à mettre
# à jour.
# ---------------------------------------------------------------------------

SLIDE1_SLOTS = {
    "titre":             3,    # "l'activité et les embauches continuent..."
    "bloc_gauche":       17,   # chapeau + commentaire CA/Eff
    "bloc_droite":       9,    # commentaire offre/demande
    "source_gauche":     28,   # "Champ : Total ..."
    "source_droite":     31,
    "chart_gauche":      6,    # CA + Effectifs
    "chart_droite":      8,    # Carnets passés
    # Libellés de graphiques dans le template — à SUPPRIMER car nos PNG
    # contiennent déjà leur propre titre (évite les titres dupliqués).
    "label_chart_gauche": 26,  # "Évolution de l'activité et des effectifs..."
    "label_chart_droite": 33,  # "Jugement sur l'état des carnets..."
}

SLIDE2_SLOTS = {
    "titre":             3,    # "La dégradation de la conjoncture..."
    "commentaire":       10,   # gros commentaire sectoriel
    "source":            14,
    "chart":             7,    # CA par secteur
    "label_chart":       13,   # "Évolution de l'activité par branche..."
}


# ---------------------------------------------------------------------------
# Helpers : navigation et manipulation de shapes
# ---------------------------------------------------------------------------

def _find_shape_by_id(slide, shape_id: int):
    """
    Retrouve un shape par son shape_id (int). Retourne None si absent.
    """
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    return None


def _strip_markdown_bold(text: str) -> str:
    """
    Supprime les marqueurs de gras markdown (**texte** et *texte*) que
    Gemini insère parfois malgré les consignes de prompt.

    Le formatting gras est géré au niveau du XML PowerPoint par notre
    code — les astérisques du LLM ne doivent pas se retrouver dans le
    rendu final.

    Gère aussi les cas asymétriques (** ouvrant sans fermant, ou l'inverse)
    en nettoyant les astérisques isolés.
    """
    # Retire les ** appariés (double astérisque)
    text = re.sub(r"\*\*([^*]+?)\*\*", r"\1", text)
    # Retire les * isolés
    text = re.sub(r"\*([^*]+?)\*", r"\1", text)
    # Retire les ** ou * orphelins qui auraient pu rester
    text = text.replace("**", "").replace("*", "")
    # Normalise les espaces multiples créés
    text = re.sub(r" {2,}", " ", text)
    return text.strip()


def _enlarge_title_box(shape, prs) -> None:
    """
    Élargit la boîte de titre pour qu'elle occupe quasiment toute la
    largeur de la slide, ET active le retour à la ligne automatique
    pour que les titres longs s'affichent sur 2 lignes dans la boîte
    au lieu de déborder.

    Motivation : dans le template Bpifrance, les boîtes de titre sont
    configurées avec :
    - `wrap='none'` (pas de retour à la ligne automatique)
    - `spAutoFit` (la boîte s'étend pour suivre le texte)
    Résultat : un titre long continue sur une ligne unique et sort du
    cadre de la slide. Le template original contourne ça avec un
    line-break manuel au milieu du titre, mais ça n'est pas généralisable
    à un titre produit par le LLM.

    Ce qu'on fait :
    1. Élargir la boîte à 32.5 cm (la boîte slide 2 n'en faisait que 29).
    2. Remplacer `wrap='none'` par `wrap='square'` (retour à la ligne
       automatique dans la largeur de la boîte).
    3. Retirer `spAutoFit` pour figer la hauteur (empêcher la boîte de
       s'étendre vers le bas sans limite).

    La largeur cible 32.5 cm est celle utilisée par Bpifrance sur la
    slide 1 — c'est la largeur maximale utile, le `left=-1.8 cm`
    légèrement négatif créant un effet "full bleed" voulu.
    """
    from lxml import etree
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    # 1. Largeur cible
    target_width_emu = int(32.5 * 360000)
    shape.width = target_width_emu

    # 2. Forcer le text wrap automatique
    txBody = shape._element.find(
        ".//{http://schemas.openxmlformats.org/drawingml/2006/main}txBody"
    )
    if txBody is None:
        # Namespace alternatif DrawingML pour les shapes (pic/sp)
        txBody = shape._element.find(
            ".//{http://schemas.openxmlformats.org/presentationml/2006/main}txBody"
        )
    if txBody is None:
        return

    bodyPr = txBody.find(f"{ns_a}bodyPr")
    if bodyPr is None:
        return

    # Activer wrap='square' (retour à la ligne automatique)
    bodyPr.set("wrap", "square")

    # 3. Retirer spAutoFit (qui agrandissait la boîte pour suivre le texte).
    # On le remplace par <a:noAutofit/> pour figer explicitement la taille.
    for child in list(bodyPr):
        tag = child.tag.split("}")[-1]
        if tag in ("spAutoFit", "normAutofit"):
            bodyPr.remove(child)
    # Ajouter noAutofit si absent (peut déjà avoir été supprimé)
    if bodyPr.find(f"{ns_a}noAutofit") is None:
        bodyPr.append(etree.SubElement(bodyPr, f"{ns_a}noAutofit"))


def _enforce_title_length(title: str, max_words: int = 18) -> str:
    """
    Garde-fou pour éviter les titres trop longs qui débordent de la slide.

    Les boîtes de titre du template Bpifrance acceptent jusqu'à ~18 mots
    sur 2 lignes (avec police 44pt), à condition d'élargir la slide 2
    comme le fait report_builder. Ce garde-fou est une sécurité de
    dernier recours si Gemini part en roue libre.

    Tronqué en préservant un point ou une virgule proche quand possible,
    pour ne pas couper au milieu d'une idée.
    """
    title = title.strip()
    # Retire la ponctuation finale (le template Bpifrance n'en met pas).
    title = title.rstrip(".!?")
    words = title.split()
    if len(words) <= max_words:
        return title
    truncated = " ".join(words[:max_words])
    print(f"  ⚠ titre tronqué ({len(words)} mots → {max_words}) : {title!r}")
    return truncated


def _clear_runs_keeping_first(paragraph) -> None:
    """Supprime tous les runs d'un paragraphe sauf le premier."""
    runs = list(paragraph.runs)
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


# Regex pour découper une phrase française : ponctuation de fin (. ! ?)
# suivie d'un espace et d'une majuscule. On capture le séparateur pour le
# conserver avec la 1ère phrase.
_SENTENCE_SPLIT_RE = re.compile(r"(?<=[.!?])\s+(?=[A-ZÀÂÉÈÊÎÔÛÇ])")


def _split_first_sentence(text: str) -> tuple[str, str]:
    """
    Sépare un texte en (première_phrase, reste).

    La première phrase inclut son signe de ponctuation final. Le 'reste'
    commence par un espace pour que la concaténation visuelle soit propre
    quand la 1ère phrase est en gras et le reste en regular.

    Exemples :
        "A. B. C."   -> ("A.", " B. C.")
        "A ! B."     -> ("A !", " B.")
        "A."         -> ("A.", "")
        "A sans point" -> ("A sans point", "")
    """
    text = text.strip()
    if not text:
        return "", ""
    match = _SENTENCE_SPLIT_RE.search(text)
    if match is None:
        # Une seule phrase (ou pas de séparateur trouvé) -> tout en gras.
        return text, ""
    first = text[:match.start()].rstrip()
    rest = " " + text[match.end():].lstrip()
    return first, rest


def _replace_paragraph_text_keep_style(paragraph, new_text: str) -> None:
    """
    Remplace le texte d'un paragraphe en gardant EXACTEMENT le style du
    1er run existant (police, couleur, gras, italique). Les autres runs
    sont supprimés. Utilisé pour les paragraphes mono-run : titres,
    chapeaux verts, lignes de source.
    """
    runs = list(paragraph.runs)
    if not runs:
        paragraph.add_run().text = new_text
        return
    runs[0].text = new_text
    _clear_runs_keeping_first(paragraph)


def _replace_paragraph_with_bold_first_sentence(paragraph, new_text: str) -> None:
    """
    Remplace le texte d'un paragraphe en créant 2 runs :
    - run 1 : première phrase, en GRAS, héritant du style du 1er run original
      (police, couleur, taille).
    - run 2 : reste du paragraphe, NON GRAS, même style pour le reste.

    Usage : paragraphes à puce "noirs" où seule la 1ère phrase est mise en
    valeur (pattern Bpifrance typique).

    Si le paragraphe du template n'a pas de run existant (rare), on
    retombe sur un simple .text = new_text.
    """
    first, rest = _split_first_sentence(new_text)

    runs = list(paragraph.runs)
    if not runs:
        # Aucun style à hériter. Fallback minimal.
        paragraph.add_run().text = new_text
        return

    # On garde le 1er run du template (qui a le style de référence) et on
    # va s'en servir pour la 1ère phrase en gras.
    first_run = runs[0]
    first_run.text = first
    first_run.font.bold = True

    # Supprimer les runs suivants du template (résidus).
    _clear_runs_keeping_first(paragraph)

    # Ajouter un 2ème run pour le reste, en copiant le <a:rPr> du 1er
    # puis en forçant bold=False. On passe par XML pour avoir un vrai
    # clone des propriétés de run (sinon on perd la couleur héritée).
    if rest:
        first_r_elem = first_run._r
        # Trouver <a:rPr> du 1er run (peut être absent si style pur hérité).
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        rPr_src = first_r_elem.find(f"{ns_a}rPr")

        # Créer un nouveau run <a:r> à côté du premier.
        new_r = deepcopy(first_r_elem)
        # Vider son texte actuel
        t_elem = new_r.find(f"{ns_a}t")
        if t_elem is not None:
            t_elem.text = rest
        # Forcer bold=0 sur le rPr cloné
        new_rPr = new_r.find(f"{ns_a}rPr")
        if new_rPr is not None:
            new_rPr.set("b", "0")
        # Insérer après le 1er run
        first_r_elem.addnext(new_r)


def _replace_text(shape, new_text: str) -> None:
    """
    Remplace le texte d'un shape en préservant le formatting XML des
    paragraphes et runs EXISTANTS (polices, couleurs, puces, niveaux,
    capitalisation, espacements, etc.).

    Stratégie :
    - Split `new_text` en paragraphes via '\\n\\n' (séparateur recommandé).
    - Le text_frame du template a N_template paragraphes, chacun avec un
      style XML propre (ex. chapeau = vert gras majuscules, corps = noir
      normal).
    - On mappe : paragraphe i du template <- texte i du nouveau contenu.
      Le formatting du paragraphe i est PRÉSERVÉ, seul le texte change.
    - Si plus de paragraphes neufs que dans le template : on clone
      (deepcopy XML) le DERNIER paragraphe du template (qui est
      généralement le style "corps").
    - Si moins : on supprime les paragraphes en trop.

    Pourquoi on modifie le 1er run en place (vs en recréer un) :
    le XML du 1er run contient toutes les propriétés run-level
    (`<a:rPr>`) — si on le remplace on perd capitalisation, soulignement,
    tailles spéciales, etc. Modifier juste son .text est non-destructif.

    Cas robustes :
    - Shape sans text_frame : warning + retour silencieux.
    - Paragraphe sans run (rare, possible sur placeholder vide) : on
      crée un run héritant du paragraphe.
    - `new_text` vide : laisse un paragraphe vide.
    """
    if not shape.has_text_frame:
        print(f"  ⚠ shape {shape.shape_id} sans text_frame, ignoré")
        return

    tf = shape.text_frame
    target_paragraphs = new_text.split("\n\n") if new_text else [""]
    n_new = len(target_paragraphs)
    existing = list(tf.paragraphs)
    n_old = len(existing)

    if n_old == 0:
        # Cas pathologique : text_frame sans paragraphe. On crée un
        # paragraphe par défaut (rare, mais safety net).
        tf.text = target_paragraphs[0]
        for para_text in target_paragraphs[1:]:
            tf.add_paragraph().text = para_text
        return

    # Étape 1 : aligner le nombre de paragraphes
    if n_new > n_old:
        # Cloner le dernier paragraphe existant pour matcher la taille
        template_p = existing[-1]._p
        for _ in range(n_new - n_old):
            cloned = deepcopy(template_p)
            template_p.addnext(cloned)
            template_p = cloned
        existing = list(tf.paragraphs)
    elif n_new < n_old:
        # Supprimer les paragraphes excédentaires
        for p in existing[n_new:]:
            p._p.getparent().remove(p._p)
        existing = existing[:n_new]

    # Étape 2 : pour chaque paragraphe, remplacer le texte du 1er run
    # (et vider les autres runs pour éviter des résidus du template).
    for p, text in zip(existing, target_paragraphs):
        runs = list(p.runs)
        if not runs:
            # Pas de run : on en crée un. Il héritera du <a:pPr> du paragraphe.
            p.add_run().text = text
        else:
            runs[0].text = text
            _clear_runs_keeping_first(p)


def _fill_commentary_block(shape,
                           chapeau: str,
                           puces: list[str]) -> None:
    """
    Remplit un shape du template en préservant la structure Bpifrance :
    - 1er paragraphe : CHAPEAU (vert gras, 1 seul run tout gras)
    - paragraphes suivants : PUCES (noires, 1ère phrase en gras + reste en regular)

    Le template fournit DÉJÀ les styles XML par paragraphe :
    - paragraphe 0 : chapeau vert (tout le run est en gras vert dans le XML)
    - paragraphes 1..N : puces à lvl=3 avec puce jaune, 1er run en gras noir

    On ne change JAMAIS les propriétés XML de paragraphe (pPr) ni les puces :
    on se contente de :
    1. Nettoyer le markdown ** / * que Gemini peut introduire malgré nos consignes.
    2. Filtrer les paragraphes "vides" du template (sans run) qui traînent
       parfois en fin de text_frame et qu'on ne peut pas cloner proprement.
    3. aligner le nombre de paragraphes-puces avec le contenu à injecter.
    4. pour le chapeau : remplacer le texte en préservant le run existant.
    5. pour chaque puce : découper en (1ère phrase, reste) et créer 2 runs
       qui héritent du style du 1er run du template.

    Args:
        shape   : le shape PowerPoint à remplir.
        chapeau : texte du chapeau (1 phrase). Si vide, le paragraphe
                  chapeau du template est supprimé.
        puces   : liste de paragraphes-puces à insérer.
    """
    if not shape.has_text_frame:
        print(f"  ⚠ shape {shape.shape_id} sans text_frame, ignoré")
        return

    # Nettoyage markdown en amont — les astérisques Gemini parasites
    # (**texte**) ne doivent jamais arriver au XML.
    chapeau = _strip_markdown_bold(chapeau) if chapeau else ""
    puces = [_strip_markdown_bold(p) for p in puces]

    tf = shape.text_frame
    existing = list(tf.paragraphs)

    # Filtrer les paragraphes "vides" du template (sans run) : ils ne
    # portent pas de style utilisable pour le clonage. Le template
    # Bpifrance en a parfois un en fin de text_frame.
    def _has_runs(p):
        return len(list(p.runs)) > 0

    # On garde les paras "utiles" (avec runs) en tête, et on supprime
    # les para-poubelle sans run.
    useful = [p for p in existing if _has_runs(p)]
    empty_tail = [p for p in existing if not _has_runs(p)]
    for p in empty_tail:
        parent = p._p.getparent()
        if parent is not None:
            parent.remove(p._p)
    existing = useful

    if len(existing) < 2:
        print(f"  ⚠ shape {shape.shape_id} : template attendu avec ≥2 paragraphes utiles, "
              f"trouvé {len(existing)}. Fallback sur _replace_text.")
        full_text = (chapeau + "\n\n" if chapeau else "") + "\n\n".join(puces)
        _replace_text(shape, full_text)
        return

    # Paragraphe 0 = chapeau du template (vert). Paragraphes 1..N = puces.
    chapeau_p = existing[0]
    puce_paragraphs = existing[1:]

    # --- Chapeau ---
    if chapeau:
        _replace_paragraph_text_keep_style(chapeau_p, chapeau)
    else:
        # Pas de chapeau demandé : on supprime le paragraphe chapeau du template.
        chapeau_p._p.getparent().remove(chapeau_p._p)

    # --- Puces ---
    n_old = len(puce_paragraphs)
    n_new = len(puces)

    if n_new > n_old:
        # Cloner la dernière puce existante (qui a ses runs stylés). On
        # clone TOUJOURS le même paragraphe de référence (pas chaque
        # nouveau clone), pour garantir que tous les ajouts ont le même
        # style de base sans cumul de modifications XML.
        reference_puce_xml = puce_paragraphs[-1]._p
        last_p = puce_paragraphs[-1]._p
        for _ in range(n_new - n_old):
            cloned = deepcopy(reference_puce_xml)
            last_p.addnext(cloned)
            last_p = cloned
        # Relire la liste après ajout
        all_paras = list(tf.paragraphs)
        start = 1 if chapeau else 0
        puce_paragraphs = all_paras[start:start + n_new]
    elif n_new < n_old:
        # Supprimer les puces en trop
        for p in puce_paragraphs[n_new:]:
            parent = p._p.getparent()
            if parent is not None:
                parent.remove(p._p)
        puce_paragraphs = puce_paragraphs[:n_new]

    # Pour chaque puce : 2 runs (gras + regular)
    for p, texte_puce in zip(puce_paragraphs, puces):
        _replace_paragraph_with_bold_first_sentence(p, texte_puce)


def _split_paragraphs(text: str) -> list[str]:
    """Split un texte en paragraphes sur '\\n\\n' et strip chaque morceau."""
    return [p.strip() for p in text.split("\n\n") if p.strip()]


def _remove_shape(shape) -> None:
    """
    Supprime un shape de son slide via XML surgery (python-pptx n'a pas
    d'API officielle pour ça, mais cette méthode est standard et robuste).

    Note : on stocke parent + sp dans des variables locales AVANT le
    remove, car lxml peut invalider certaines références intermédiaires
    entre deux appels .getparent() consécutifs.
    """
    sp = shape._element
    parent = sp.getparent()
    if parent is None:
        # Le shape a déjà été retiré ou n'a pas de parent (orphelin).
        return
    parent.remove(sp)


def _replace_chart_with_image(slide, chart_shape_id: int,
                              image_path: str | Path) -> None:
    """
    Remplace un CHART natif par une image PNG à la même position et
    dimensions.

    Étapes :
    1. Localise le shape chart par son ID.
    2. Capture ses coordonnées (left, top, width, height) en EMU.
    3. Le supprime.
    4. Insère l'image aux mêmes coordonnées.
    """
    chart_shape = _find_shape_by_id(slide, chart_shape_id)
    if chart_shape is None:
        raise ValueError(f"Chart shape_id={chart_shape_id} introuvable.")

    left, top = chart_shape.left, chart_shape.top
    width, height = chart_shape.width, chart_shape.height

    _remove_shape(chart_shape)
    slide.shapes.add_picture(str(image_path), left, top, width, height)


# ---------------------------------------------------------------------------
# Constructeur principal
# ---------------------------------------------------------------------------

def _format_source(n_repondants: Optional[int] = None) -> str:
    """Mention source Bpifrance standardisée, identique à celle de style.py."""
    champ = f"Total (N = {n_repondants})" if n_repondants else "Total"
    return f"Champ : {champ} ; Source : Bpifrance Le Lab"


def build_report(
    report: FullReport,
    commentary: FullCommentary,
    charts: dict[str, Path],
    template_path: str | Path = "presentation.pptx",
    output_path: str | Path = "output/rapport_conjoncture.pptx",
    n_repondants: Optional[int] = None,
) -> Path:
    """
    Construit le rapport PPTX final.

    Args:
        report       : FullReport de la phase 1 (métriques calculées).
        commentary   : FullCommentary de la phase 3 (textes générés).
        charts       : dict {"ca_eff": Path, "carnets": Path, "secteurs": Path, ...}
                       produit par charts.render_all_charts().
        template_path: PPTX de référence Bpifrance à utiliser comme template.
        output_path  : destination du rapport final.
        n_repondants : N de l'enquête (ex. 4722 pour S1 2026). Si fourni,
                       met à jour les lignes "Champ : Total (N = ...)"

    Returns:
        Path vers le fichier PPTX généré.
    """
    template_path = Path(template_path)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if not template_path.exists():
        raise FileNotFoundError(f"Template introuvable : {template_path}")

    prs = Presentation(str(template_path))
    source_text = _format_source(n_repondants) if n_repondants else None

    # -------------------- Slide 1 --------------------
    print("  Remplissage slide 1 ...")
    slide1 = prs.slides[0]

    # Titre (1 seul run mono-style : replacement simple)
    # Garde-fou de longueur : si Gemini ignore la consigne "max 18 mots",
    # on tronque côté code pour éviter le débordement visuel.
    shape = _find_shape_by_id(slide1, SLIDE1_SLOTS["titre"])
    # Normaliser la largeur de la boîte de titre (cf. _enlarge_title_box).
    _enlarge_title_box(shape, prs)
    titre_slide1 = _enforce_title_length(
        _strip_markdown_bold(commentary.slide1_title))
    _replace_text(shape, titre_slide1)

    # Bloc gauche : chapeau vert + puces CA/Effectifs (1ère phrase en gras)
    shape = _find_shape_by_id(slide1, SLIDE1_SLOTS["bloc_gauche"])
    puces_gauche = _split_paragraphs(commentary.ca_eff_comment)
    _fill_commentary_block(shape,
                           chapeau=commentary.slide1_chapeau,
                           puces=puces_gauche)

    # Bloc droite : chapeau vert + puces offre/demande
    shape = _find_shape_by_id(slide1, SLIDE1_SLOTS["bloc_droite"])
    puces_droite = _split_paragraphs(commentary.carnets_appro_comment)
    _fill_commentary_block(shape,
                           chapeau=commentary.slide1_chapeau_droite,
                           puces=puces_droite)

    # Sources (si N fourni)
    if source_text:
        for slot in ("source_gauche", "source_droite"):
            shape = _find_shape_by_id(slide1, SLIDE1_SLOTS[slot])
            if shape:
                _replace_text(shape, source_text)

    # Charts : suppression + insertion PNG
    # On retire d'abord les libellés de graphique du template (évite
    # la duplication avec ceux intégrés dans nos PNG).
    for slot in ("label_chart_gauche", "label_chart_droite"):
        shape = _find_shape_by_id(slide1, SLIDE1_SLOTS[slot])
        if shape:
            _remove_shape(shape)
    _replace_chart_with_image(slide1, SLIDE1_SLOTS["chart_gauche"],
                              charts["ca_eff"])
    _replace_chart_with_image(slide1, SLIDE1_SLOTS["chart_droite"],
                              charts["carnets"])

    # -------------------- Slide 2 --------------------
    print("  Remplissage slide 2 ...")
    slide2 = prs.slides[1]

    shape = _find_shape_by_id(slide2, SLIDE2_SLOTS["titre"])
    # Élargir la boîte de titre (de 29 à 32.5 cm) pour que les titres
    # de 14-18 mots tiennent sur 2 lignes au lieu de 3.
    _enlarge_title_box(shape, prs)
    titre_slide2 = _enforce_title_length(
        _strip_markdown_bold(commentary.slide2_title))
    _replace_text(shape, titre_slide2)

    # Bloc commentaire sectoriel : puce verte d'intro (chapeau) + puces détaillées
    # IMPORTANT : sur la slide 2, le "paragraphe 0" du template est DÉJÀ
    # une puce verte (lvl=3, couleur verte). On l'utilise pour le chapeau
    # synthétique ; les paragraphes suivants sont les puces noires détaillées.
    shape = _find_shape_by_id(slide2, SLIDE2_SLOTS["commentaire"])
    puces_sect = _split_paragraphs(commentary.secteurs_comment)
    _fill_commentary_block(shape,
                           chapeau=commentary.slide2_chapeau,
                           puces=puces_sect)

    if source_text:
        shape = _find_shape_by_id(slide2, SLIDE2_SLOTS["source"])
        if shape:
            _replace_text(shape, source_text)

    # Suppression du libellé de graphique du template (doublon avec PNG)
    shape = _find_shape_by_id(slide2, SLIDE2_SLOTS["label_chart"])
    if shape:
        _remove_shape(shape)

    # ALIGNEMENT VISUEL : sur la slide 2, le template met le graphique
    # (shape 7, top=5.7cm) un peu plus bas que le commentaire (shape 10,
    # top=4.3cm). Ce décalage était visuellement "comblé" par le libellé
    # du graphique (shape 13) qu'on vient de supprimer. Sans ce libellé,
    # on veut que le graphique commence au même niveau vertical que le
    # commentaire à gauche. On lit le top du commentaire et on aligne.
    commentaire_shape = _find_shape_by_id(slide2, SLIDE2_SLOTS["commentaire"])
    if commentaire_shape is not None:
        chart_shape = _find_shape_by_id(slide2, SLIDE2_SLOTS["chart"])
        if chart_shape is not None:
            # On aligne et on étend la hauteur disponible pour
            # que le graphique reste visible sans chevaucher la source.
            target_top = commentaire_shape.top
            old_top = chart_shape.top
            old_height = chart_shape.height
            # Nouveau haut = top du commentaire
            chart_shape.top = target_top
            # On garde la même hauteur qu'avant (le graphique est au
            # ratio 16:9, on n'a pas besoin de l'étirer).
            chart_shape.height = old_height

    _replace_chart_with_image(slide2, SLIDE2_SLOTS["chart"],
                              charts["secteurs"])

    # -------------------- Sauvegarde --------------------
    prs.save(str(output_path))
    print(f"  ✓ rapport sauvegardé : {output_path}")
    return output_path
